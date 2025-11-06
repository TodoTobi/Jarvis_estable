# actions_windows.py
import os
import subprocess
import webbrowser
import shutil
import json
from pathlib import Path
from datetime import datetime
from typing import List, Optional
from urllib.parse import quote

# Intentar importar librerías opcionales
try:
    import pyautogui
    PYAutoGUI_AVAILABLE = True
except ImportError:
    PYAutoGUI_AVAILABLE = False

try:
    from PIL import ImageGrab
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

# carpeta del usuario actual
USER_HOME = os.path.expanduser("~")
DESKTOP = os.path.join(USER_HOME, "Desktop")
JARVIS_TRASH = os.path.join(USER_HOME, ".jarvis", "trash")

# Crear carpeta de papelera si no existe
os.makedirs(JARVIS_TRASH, exist_ok=True)

# ========================
# OPERACIONES DE ARCHIVOS Y CARPETAS (BÁSICAS)
# ========================

def crear_carpeta(ruta: str) -> str:
    """Crea una carpeta en la ruta indicada."""
    os.makedirs(ruta, exist_ok=True)
    return ruta

def crear_carpeta_en_escritorio(nombre: str) -> str:
    """Crea una carpeta en el escritorio."""
    ruta = os.path.join(DESKTOP, nombre)
    os.makedirs(ruta, exist_ok=True)
    return ruta

def listar_carpeta(ruta: str, filtro: Optional[str] = None) -> List[dict]:
    """Lista archivos y subcarpetas. Opcionalmente filtra por extensión."""
    if not os.path.exists(ruta):
        raise FileNotFoundError(f"No existe la ruta: {ruta}")
    
    items = []
    try:
        for entry in os.scandir(ruta):
            item = {
                "nombre": entry.name,
                "es_carpeta": entry.is_dir(),
                "tamaño": entry.stat().st_size if entry.is_file() else 0
            }
            if filtro:
                if entry.is_file() and entry.name.endswith(filtro):
                    items.append(item)
                elif not filtro.startswith("."):
                    items.append(item)
            else:
                items.append(item)
    except PermissionError:
        raise PermissionError(f"No tienes permisos para acceder a: {ruta}")
    
    return items

def leer_archivo(ruta: str, limite_mb: int = 5) -> str:
    """Lee un archivo de texto con límite de tamaño."""
    if not os.path.exists(ruta):
        raise FileNotFoundError(f"No existe el archivo: {ruta}")
    
    # Verificar tamaño
    tamaño = os.path.getsize(ruta)
    if tamaño > limite_mb * 1024 * 1024:
        raise ValueError(f"Archivo muy grande ({tamaño / 1024 / 1024:.2f}MB). Límite: {limite_mb}MB")
    
    try:
        with open(ruta, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        raise ValueError("El archivo no es de texto o tiene codificación no UTF-8")

def crear_txt(ruta: str, contenido: str = "") -> str:
    """Crea o sobrescribe un archivo de texto."""
    carpeta = os.path.dirname(ruta)
    if carpeta and not os.path.exists(carpeta):
        os.makedirs(carpeta, exist_ok=True)
    with open(ruta, "w", encoding="utf-8") as f:
        f.write(contenido)
    return ruta

def editar_archivo(ruta: str, contenido: str, modo: str = "sobrescribir") -> str:
    """Edita un archivo. Modos: 'sobrescribir', 'agregar'."""
    if modo == "agregar":
        with open(ruta, "a", encoding="utf-8") as f:
            f.write(contenido)
    else:
        with open(ruta, "w", encoding="utf-8") as f:
            f.write(contenido)
    return ruta

def copiar_archivo(origen: str, destino: str) -> str:
    """Copia un archivo o carpeta manteniendo permisos."""
    if os.path.isdir(origen):
        shutil.copytree(origen, destino, dirs_exist_ok=True)
    else:
        shutil.copy2(origen, destino)
    return destino

def mover_archivo(origen: str, destino: str) -> str:
    """Mueve un archivo o carpeta."""
    shutil.move(origen, destino)
    return destino

def eliminar(ruta: str, permanente: bool = False) -> str:
    """Elimina un archivo o carpeta. Si no es permanente, lo mueve a papelera."""
    if not os.path.exists(ruta):
        raise FileNotFoundError(f"No existe: {ruta}")
    
    if permanente:
        if os.path.isdir(ruta):
            shutil.rmtree(ruta)
        else:
            os.remove(ruta)
        return f"Eliminado permanentemente: {ruta}"
    else:
        # Mover a papelera
        nombre = os.path.basename(ruta)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        destino = os.path.join(JARVIS_TRASH, f"{timestamp}_{nombre}")
        shutil.move(ruta, destino)
        return f"Movido a papelera: {destino}"

def duplicar(ruta: str, nuevo_nombre: Optional[str] = None) -> str:
    """Duplica un archivo o carpeta con nuevo nombre."""
    if not os.path.exists(ruta):
        raise FileNotFoundError(f"No existe: {ruta}")
    
    carpeta = os.path.dirname(ruta)
    if not nuevo_nombre:
        nombre_base, extension = os.path.splitext(os.path.basename(ruta))
        nuevo_nombre = f"{nombre_base}_copia{extension}"
    
    destino = os.path.join(carpeta, nuevo_nombre)
    return copiar_archivo(ruta, destino)

# ========================
# OPERACIONES DEL SISTEMA / APPS
# ========================

def abrir_carpeta(ruta: str):
    """Abre una carpeta en el Explorador de Windows."""
    if os.path.isdir(ruta):
        os.startfile(ruta)
    else:
        raise FileNotFoundError(f"No existe la carpeta: {ruta}")

def encontrar_navegador(nombre: str) -> Optional[str]:
    """Encuentra el path completo de un navegador."""
    nombre_lower = nombre.lower()
    
    # Rutas comunes de navegadores en Windows
    rutas_comunes = [
        os.path.expanduser(r"~\AppData\Local\Google\Chrome\Application\chrome.exe"),
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
        os.path.expanduser(r"~\AppData\Local\BraveSoftware\Brave-Browser\Application\brave.exe"),
        r"C:\Program Files\BraveSoftware\Brave-Browser\Application\brave.exe",
        r"C:\Program Files (x86)\BraveSoftware\Brave-Browser\Application\brave.exe",
        r"C:\Program Files\Mozilla Firefox\firefox.exe",
        r"C:\Program Files (x86)\Mozilla Firefox\firefox.exe",
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    ]
    
    # Buscar Chrome
    if "chrome" in nombre_lower or "google chrome" in nombre_lower:
        for ruta in rutas_comunes:
            if "chrome.exe" in ruta.lower() and os.path.exists(ruta):
                return ruta
        # Intentar con el nombre simple
        try:
            subprocess.run(["where", "chrome.exe"], capture_output=True, check=True)
            return "chrome.exe"
        except:
            pass
    
    # Buscar Brave
    if "brave" in nombre_lower:
        for ruta in rutas_comunes:
            if "brave" in ruta.lower() and os.path.exists(ruta):
                return ruta
    
    # Buscar Firefox
    if "firefox" in nombre_lower:
        for ruta in rutas_comunes:
            if "firefox.exe" in ruta.lower() and os.path.exists(ruta):
                return ruta
    
    # Buscar Edge
    if "edge" in nombre_lower:
        for ruta in rutas_comunes:
            if "msedge.exe" in ruta.lower() and os.path.exists(ruta):
                return ruta
    
    return None

def abrir_app(nombre: str):
    """Ejecuta un programa por nombre o path."""
    nombre_lower = nombre.lower()
    path_app = None
    
    # Si es un navegador, buscar su path
    if any(nav in nombre_lower for nav in ["chrome", "brave", "firefox", "edge", "navegador"]):
        path_app = encontrar_navegador(nombre)
        if path_app:
            subprocess.Popen([path_app], shell=False)
            return
    
    # Normalizar nombre de aplicación común
    if "chrome" in nombre_lower and not nombre_lower.endswith(".exe"):
        nombre = "chrome.exe"
    elif "firefox" in nombre_lower and not nombre_lower.endswith(".exe"):
        nombre = "firefox.exe"
    elif "edge" in nombre_lower and not nombre_lower.endswith(".exe"):
        nombre = "msedge.exe"
    elif "brave" in nombre_lower:
        nombre = "brave.exe"
    
    # Intentar ejecutar directamente
    try:
        subprocess.Popen(nombre, shell=True)
    except:
        # Si falla, intentar encontrar el ejecutable
        try:
            result = subprocess.run(["where", nombre], capture_output=True, text=True, check=True)
            path = result.stdout.strip().split('\n')[0]
            subprocess.Popen([path], shell=False)
        except:
            # Último intento: usar start en Windows
            subprocess.Popen(f'start "" "{nombre}"', shell=True)

def cerrar_app_por_nombre(nombre: str):
    """Cierra una aplicación por nombre de proceso."""
    subprocess.run(["taskkill", "/f", "/im", nombre], 
                   stdout=subprocess.DEVNULL, 
                   stderr=subprocess.DEVNULL)

def cerrar_ventana_activa():
    """Cierra la ventana activa (Alt+F4)."""
    if PYAutoGUI_AVAILABLE:
        pyautogui.hotkey("alt", "f4")
    else:
        # Fallback: intentar cerrar proceso foreground (menos preciso)
        subprocess.run(["taskkill", "/f", "/fi", "status eq running"],
                       stdout=subprocess.DEVNULL, 
                       stderr=subprocess.DEVNULL)

def tomar_screenshot(ruta: Optional[str] = None) -> str:
    """Toma un screenshot y lo guarda."""
    if not PIL_AVAILABLE:
        raise ImportError("PIL/Pillow no está instalado. Instalá: pip install Pillow")
    
    if not ruta:
        ruta = os.path.join(DESKTOP, f"screenshot_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png")
    
    screenshot = ImageGrab.grab()
    screenshot.save(ruta)
    return ruta

# ========================
# NAVEGADOR / WEB
# ========================

def abrir_url(url: str, nueva_pestaña: bool = True, navegador: Optional[str] = None):
    """Abre una URL en el navegador."""
    if navegador:
        # Abrir en navegador específico
        path_nav = encontrar_navegador(navegador)
        if path_nav:
            subprocess.Popen([path_nav, url], shell=False)
        else:
            # Fallback: usar webbrowser con el navegador por defecto
            webbrowser.open(url)
    else:
        if nueva_pestaña:
            webbrowser.open_new_tab(url)
        else:
            webbrowser.open(url)

def buscar_google(consulta: str):
    """Busca en Google."""
    url = f"https://www.google.com/search?q={consulta.replace(' ', '+')}"
    webbrowser.open_new_tab(url)

def buscar_youtube(consulta: str, navegador: Optional[str] = None):
    """Busca en YouTube y abre la URL con la búsqueda."""
    consulta_encoded = quote(consulta)
    url = f"https://www.youtube.com/results?search_query={consulta_encoded}"
    
    if navegador:
        path_nav = encontrar_navegador(navegador)
        if path_nav:
            subprocess.Popen([path_nav, url], shell=False)
        else:
            webbrowser.open_new_tab(url)
    else:
        webbrowser.open_new_tab(url)

def abrir_youtube_en_navegador(navegador: str):
    """Abre YouTube en un navegador específico."""
    url = "https://www.youtube.com"
    path_nav = encontrar_navegador(navegador)
    if path_nav:
        subprocess.Popen([path_nav, url], shell=False)
    else:
        webbrowser.open(url)

# ========================
# CONTROL DE PLATAFORMAS (YouTube, TikTok, Instagram)
# ========================

def control_youtube(accion: str, navegador: Optional[str] = None):
    """Controla YouTube: pausar, reproducir, siguiente, anterior, volumen, lista."""
    # Usar pyautogui para controlar YouTube
    if not PYAutoGUI_AVAILABLE:
        raise ImportError("pyautogui no está instalado. Instalá: pip install pyautogui")
    
    # Abrir YouTube si no está especificado el navegador
    if navegador:
        path_nav = encontrar_navegador(navegador)
        if path_nav:
            subprocess.Popen([path_nav, "https://www.youtube.com"], shell=False)
            import time
            time.sleep(2)  # Esperar a que cargue
    
    accion_lower = accion.lower()
    
    # Mapeo de acciones a atajos de teclado
    if "pausar" in accion_lower or "pause" in accion_lower:
        pyautogui.press('space')
    elif "reproducir" in accion_lower or "play" in accion_lower:
        pyautogui.press('space')
    elif "siguiente" in accion_lower or "next" in accion_lower:
        pyautogui.press('shift+n')  # Siguiente video
    elif "anterior" in accion_lower or "previous" in accion_lower or "atras" in accion_lower:
        pyautogui.press('shift+p')  # Video anterior
    elif "lista" in accion_lower or "playlist" in accion_lower:
        pyautogui.press('shift+l')  # Mostrar/ocultar lista
    elif "volumen arriba" in accion_lower or "subir volumen" in accion_lower or "volume up" in accion_lower:
        pyautogui.press('up', presses=5)  # Subir volumen
    elif "volumen abajo" in accion_lower or "bajar volumen" in accion_lower or "volume down" in accion_lower:
        pyautogui.press('down', presses=5)  # Bajar volumen
    elif "mute" in accion_lower or "silenciar" in accion_lower:
        pyautogui.press('m')  # Mute/unmute
    else:
        raise ValueError(f"Acción de YouTube no reconocida: {accion}")

def control_tiktok(accion: str, navegador: Optional[str] = None):
    """Controla TikTok: pausar, siguiente, anterior, like, etc."""
    if not PYAutoGUI_AVAILABLE:
        raise ImportError("pyautogui no está instalado.")
    
    if navegador:
        path_nav = encontrar_navegador(navegador)
        if path_nav:
            subprocess.Popen([path_nav, "https://www.tiktok.com"], shell=False)
            import time
            time.sleep(2)
    
    accion_lower = accion.lower()
    
    if "pausar" in accion_lower or "pause" in accion_lower:
        pyautogui.click()  # Click en el video pausa/reproduce
    elif "siguiente" in accion_lower or "next" in accion_lower:
        pyautogui.press('down')  # Scroll down = siguiente
    elif "anterior" in accion_lower or "previous" in accion_lower:
        pyautogui.press('up')  # Scroll up = anterior
    elif "like" in accion_lower:
        pyautogui.doubleClick()  # Doble click = like
    else:
        raise ValueError(f"Acción de TikTok no reconocida: {accion}")

def control_instagram(accion: str, navegador: Optional[str] = None):
    """Controla Instagram: siguiente historia, like, comentar, etc."""
    if not PYAutoGUI_AVAILABLE:
        raise ImportError("pyautogui no está instalado.")
    
    if navegador:
        path_nav = encontrar_navegador(navegador)
        if path_nav:
            subprocess.Popen([path_nav, "https://www.instagram.com"], shell=False)
            import time
            time.sleep(2)
    
    accion_lower = accion.lower()
    
    if "siguiente" in accion_lower or "next" in accion_lower:
        pyautogui.press('right')  # Siguiente historia/post
    elif "anterior" in accion_lower or "previous" in accion_lower:
        pyautogui.press('left')  # Anterior
    elif "like" in accion_lower:
        pyautogui.press('l')  # Like
    elif "comentar" in accion_lower:
        pyautogui.press('c')  # Comentar
    else:
        raise ValueError(f"Acción de Instagram no reconocida: {accion}")

# ========================
# GOOGLE DOCS Y OFFICE
# ========================

def crear_doc_google_docs(nombre: str, plantilla: Optional[str] = None, contenido: Optional[str] = None):
    """Crea un documento en Google Docs. Si hay plantilla, la duplica."""
    # Abrir Google Docs
    if plantilla:
        # Buscar plantilla y abrirla para duplicar
        url = f"https://docs.google.com/document/create?name={quote(nombre)}"
        # Nota: Para duplicar una plantilla específica, necesitarías el ID del documento
        # Por ahora abrimos Docs y el usuario puede duplicar manualmente
        webbrowser.open(url)
        return f"Documento '{nombre}' creado. Para usar plantilla '{plantilla}', duplicala manualmente."
    else:
        url = f"https://docs.google.com/document/create?name={quote(nombre)}"
        webbrowser.open(url)
        return f"Documento '{nombre}' creado en Google Docs"

def crear_docx(ruta: str, contenido: str = "", plantilla: Optional[str] = None):
    """Crea un archivo .docx. Si hay plantilla, la copia."""
    try:
        from docx import Document
        from docx.shared import Inches
    except ImportError:
        raise ImportError("python-docx no está instalado. Instalá: pip install python-docx")
    
    if plantilla and os.path.exists(plantilla):
        # Copiar plantilla
        shutil.copy2(plantilla, ruta)
        doc = Document(ruta)
    else:
        doc = Document()
    
    if contenido:
        # Agregar contenido
        for linea in contenido.split('\n'):
            if linea.strip():
                doc.add_paragraph(linea)
    
    doc.save(ruta)
    return ruta

def crear_ppt(ruta: str, titulo: str = "Presentación", plantilla: Optional[str] = None):
    """Crea un archivo .pptx. Si hay plantilla, la copia."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx no está instalado. Instalá: pip install python-pptx")
    
    if plantilla and os.path.exists(plantilla):
        prs = Presentation(plantilla)
    else:
        prs = Presentation()
        # Agregar slide de título
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = titulo
    
    prs.save(ruta)
    return ruta

# ========================
# INTEGRACIONES (Canva, Gamma, Spotify)
# ========================

def abrir_canva():
    """Abre Canva en el navegador."""
    webbrowser.open("https://www.canva.com")

def abrir_gamma():
    """Abre Gamma.app en el navegador."""
    webbrowser.open("https://gamma.app")

def control_spotify(accion: str):
    """Controla Spotify: play, pause, siguiente, anterior, volumen."""
    if not PYAutoGUI_AVAILABLE:
        raise ImportError("pyautogui no está instalado.")
    
    accion_lower = accion.lower()
    
    # Atajos de teclado globales de Spotify (si está en foco)
    if "pausar" in accion_lower or "pause" in accion_lower:
        pyautogui.press('space')
    elif "reproducir" in accion_lower or "play" in accion_lower:
        pyautogui.press('space')
    elif "siguiente" in accion_lower or "next" in accion_lower:
        pyautogui.press('next')  # Media Next
    elif "anterior" in accion_lower or "previous" in accion_lower:
        pyautogui.press('prev')  # Media Previous
    elif "volumen arriba" in accion_lower or "subir volumen" in accion_lower:
        pyautogui.press('volumeup', presses=5)
    elif "volumen abajo" in accion_lower or "bajar volumen" in accion_lower:
        pyautogui.press('volumedown', presses=5)
    else:
        # Abrir Spotify
        subprocess.Popen("spotify", shell=True)
        import time
        time.sleep(2)
        if "buscar" in accion_lower:
            # Buscar en Spotify (requiere UI automation más compleja)
            raise NotImplementedError("Búsqueda en Spotify requiere más automatización")

# ========================
# INTEGRACIÓN CON IA (ChatGPT, Gemini, Cursor)
# ========================

def enviar_prompt_chatgpt(prompt: str, navegador: Optional[str] = None):
    """Envía un prompt a ChatGPT."""
    # Abrir ChatGPT y escribir el prompt (requiere UI automation)
    if navegador:
        path_nav = encontrar_navegador(navegador)
        if path_nav:
            subprocess.Popen([path_nav, "https://chat.openai.com"], shell=False)
            import time
            time.sleep(3)
            if PYAutoGUI_AVAILABLE:
                # Escribir el prompt (requiere encontrar el textarea)
                pyautogui.write(prompt, interval=0.05)
                pyautogui.press('enter')
        else:
            webbrowser.open("https://chat.openai.com")
    else:
        webbrowser.open("https://chat.openai.com")

def enviar_prompt_gemini(prompt: str, navegador: Optional[str] = None):
    """Envía un prompt a Gemini."""
    if navegador:
        path_nav = encontrar_navegador(navegador)
        if path_nav:
            subprocess.Popen([path_nav, "https://gemini.google.com"], shell=False)
            import time
            time.sleep(3)
            if PYAutoGUI_AVAILABLE:
                pyautogui.write(prompt, interval=0.05)
                pyautogui.press('enter')
        else:
            webbrowser.open("https://gemini.google.com")
    else:
        webbrowser.open("https://gemini.google.com")

def enviar_prompt_cursor(prompt: str):
    """Envía un prompt a Cursor (si está abierto)."""
    if not PYAutoGUI_AVAILABLE:
        raise ImportError("pyautogui no está instalado.")
    
    # Abrir Cursor si no está abierto
    try:
        subprocess.Popen(["cursor"], shell=True)
        import time
        time.sleep(2)
    except:
        pass
    
    # Enviar prompt a Cursor (Ctrl+L para abrir chat)
    pyautogui.hotkey('ctrl', 'l')
    import time
    time.sleep(0.5)
    pyautogui.write(prompt, interval=0.05)
    pyautogui.press('enter')

# ========================
# ARCHIVOS AVANZADOS / ANÁLISIS
# ========================

def buscar_texto_en_archivos(ruta: str, texto: str, extensiones: Optional[List[str]] = None) -> List[dict]:
    """Busca texto dentro de archivos (grep local)."""
    resultados = []
    extensiones = extensiones or [".txt", ".py", ".js", ".html", ".md", ".json"]
    
    path = Path(ruta)
    if not path.exists():
        raise FileNotFoundError(f"No existe la ruta: {ruta}")
    
    for archivo in path.rglob("*"):
        if archivo.is_file():
            ext = archivo.suffix.lower()
            if extensiones and ext not in extensiones:
                continue
            
            try:
                with open(archivo, "r", encoding="utf-8", errors="ignore") as f:
                    contenido = f.read()
                    if texto.lower() in contenido.lower():
                        lineas = contenido.split("\n")
                        matches = [i+1 for i, linea in enumerate(lineas) if texto.lower() in linea.lower()]
                        resultados.append({
                            "archivo": str(archivo),
                            "lineas": matches[:10]  # Limitar a 10 coincidencias
                        })
            except Exception:
                continue
    
    return resultados

# ========================
# DESARROLLO / CODING AUTOMATION
# ========================

def crear_proyecto(nombre: str, template: str = "basico", ruta: Optional[str] = None) -> str:
    """Crea estructura de proyecto con carpetas y archivos base."""
    if not ruta:
        ruta = os.path.join(DESKTOP, nombre)
    
    os.makedirs(ruta, exist_ok=True)
    
    # Estructura básica
    if template == "react":
        os.makedirs(os.path.join(ruta, "src"), exist_ok=True)
        os.makedirs(os.path.join(ruta, "public"), exist_ok=True)
        crear_txt(os.path.join(ruta, "package.json"), '{"name": "' + nombre + '", "version": "1.0.0"}')
        crear_txt(os.path.join(ruta, "README.md"), f"# {nombre}\n\nProyecto creado por Jarvis.")
    elif template == "python":
        os.makedirs(os.path.join(ruta, "src"), exist_ok=True)
        crear_txt(os.path.join(ruta, "requirements.txt"), "")
        crear_txt(os.path.join(ruta, "README.md"), f"# {nombre}\n\nProyecto creado por Jarvis.")
        crear_txt(os.path.join(ruta, ".gitignore"), "__pycache__/\n*.pyc\n.env\n")
    else:  # básico
        crear_txt(os.path.join(ruta, "README.md"), f"# {nombre}\n\nProyecto creado por Jarvis.")
    
    return ruta

def abrir_vscode(ruta: Optional[str] = None):
    """Abre VSCode en la ruta especificada o directorio actual."""
    if ruta:
        subprocess.Popen(["code", ruta])
    else:
        subprocess.Popen(["code", "."])

def ejecutar_comando(comando: List[str], directorio: Optional[str] = None) -> str:
    """Ejecuta un comando del sistema."""
    result = subprocess.run(comando, cwd=directorio, capture_output=True, text=True)
    return result.stdout + result.stderr

# ========================
# UTILIDADES Y LOGS
# ========================

def guardar_log(accion: str, params: dict, resultado: str):
    """Guarda un log de acción ejecutada."""
    log_dir = os.path.join(USER_HOME, ".jarvis", "logs")
    os.makedirs(log_dir, exist_ok=True)
    
    log_file = os.path.join(log_dir, f"actions_{datetime.now().strftime('%Y%m%d')}.json")
    
    log_entry = {
        "timestamp": datetime.now().isoformat(),
        "accion": accion,
        "params": params,
        "resultado": resultado
    }
    
    logs = []
    if os.path.exists(log_file):
        try:
            with open(log_file, "r", encoding="utf-8") as f:
                logs = json.load(f)
        except:
            pass
    
    logs.append(log_entry)
    
    with open(log_file, "w", encoding="utf-8") as f:
        json.dump(logs, f, indent=2, ensure_ascii=False)

def restaurar_desde_papelera(nombre_archivo: str, destino: Optional[str] = None) -> str:
    """Restaura un archivo desde la papelera."""
    archivos = [f for f in os.listdir(JARVIS_TRASH) if f.endswith(f"_{nombre_archivo}")]
    if not archivos:
        raise FileNotFoundError(f"No se encontró '{nombre_archivo}' en la papelera")
    
    # Tomar el más reciente
    archivo = sorted(archivos)[-1]
    origen = os.path.join(JARVIS_TRASH, archivo)
    
    if not destino:
        destino = os.path.join(DESKTOP, nombre_archivo)
    
    shutil.move(origen, destino) 
    return destino

def listar_papelera() -> List[dict]:
    """Lista archivos en la papelera."""
    items = []
    for entry in os.scandir(JARVIS_TRASH):
        if entry.is_file():
            items.append({
                "nombre": entry.name,
                "fecha": datetime.fromtimestamp(entry.stat().st_mtime).isoformat(),
                "tamaño": entry.stat().st_size
            })
    return items
